"""OOXML composer that maps XML guidance to PowerPoint presentations."""

from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from xml_lib.telemetry import TelemetrySink


@dataclass
class RenderResult:
    """Result of PowerPoint rendering."""
    success: bool
    slide_count: int = 0
    citation_count: int = 0
    error: Optional[str] = None


class PPTXComposer:
    """Composes PowerPoint presentations from XML documents."""

    def __init__(
        self,
        template: Optional[Path] = None,
        telemetry: Optional[TelemetrySink] = None,
    ):
        self.template = template
        self.telemetry = telemetry

    def render(self, xml_file: Path, output_file: Path) -> RenderResult:
        """Render XML to PowerPoint presentation.

        Args:
            xml_file: Input XML file
            output_file: Output .pptx file

        Returns:
            RenderResult
        """
        try:
            # Load XML
            doc = etree.parse(str(xml_file))
            root = doc.getroot()

            # Create presentation from template or blank
            if self.template and self.template.exists():
                prs = Presentation(str(self.template))
            else:
                prs = Presentation()

            # Track citations
            citations: List[str] = []

            # Render based on root element type
            if root.tag == "document":
                self._render_document(root, prs, citations)
            else:
                # Generic rendering
                self._render_generic(root, prs, citations)

            # Add citations slide if any
            if citations:
                self._add_citations_slide(prs, citations)

            # Save presentation
            output_file.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_file))

            return RenderResult(
                success=True,
                slide_count=len(prs.slides),
                citation_count=len(citations),
            )

        except Exception as e:
            return RenderResult(
                success=False,
                error=str(e),
            )

    def _render_document(
        self,
        root: etree._Element,
        prs: Presentation,
        citations: List[str],
    ) -> None:
        """Render a document element."""
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]

        # Get title from meta
        meta = root.find("meta")
        if meta is not None:
            title_elem = meta.find("title")
            desc_elem = meta.find("description")

            if title_elem is not None and title_elem.text:
                title.text = title_elem.text
            else:
                title.text = "XML Lifecycle Document"

            if desc_elem is not None and desc_elem.text:
                subtitle.text = desc_elem.text

        # Add timestamp citation
        if root.get("timestamp"):
            citations.append(f"Document timestamp: {root.get('timestamp')}")

        # Render phases
        phases = root.find("phases")
        if phases is not None:
            for phase in phases.findall("phase"):
                self._render_phase(phase, prs, citations)

        # Summary slide
        summary = root.find("summary")
        if summary is not None:
            self._render_summary(summary, prs)

    def _render_phase(
        self,
        phase: etree._Element,
        prs: Presentation,
        citations: List[str],
    ) -> None:
        """Render a phase element as a slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        title = slide.shapes.title
        body = slide.placeholders[1]

        # Set title
        phase_name = phase.get("name", "Phase")
        title.text = f"Phase: {phase_name.title()}"

        # Add timestamp citation
        if phase.get("timestamp"):
            citations.append(f"Phase '{phase_name}' timestamp: {phase.get('timestamp')}")

        # Build content
        content_lines = []

        # Use element
        use_elem = phase.find("use")
        if use_elem is not None:
            path = use_elem.get("path", "")
            content_lines.append(f"Template: {path}")
            if use_elem.text:
                content_lines.append(f"  {use_elem.text.strip()}")
            citations.append(f"Phase '{phase_name}' references: {path}")

        # Payload
        payload = phase.find("payload")
        if payload is not None:
            content_lines.append("")
            content_lines.append("Payload:")
            for child in payload:
                text = self._element_to_text(child)
                if text:
                    content_lines.append(f"  • {text}")

        # Set body text
        if content_lines:
            text_frame = body.text_frame
            text_frame.clear()
            for i, line in enumerate(content_lines):
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                p.text = line
                p.font.size = Pt(14)

    def _render_summary(self, summary: etree._Element, prs: Presentation) -> None:
        """Render summary as final slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = "Summary"

        # Build content
        content_lines = []
        status = summary.find("status")
        if status is not None and status.text:
            content_lines.append(f"Status: {status.text}")

        next_action = summary.find("next-action")
        if next_action is not None and next_action.text:
            content_lines.append("")
            content_lines.append(f"Next Action:")
            content_lines.append(f"  {next_action.text}")

        # Set body text
        if content_lines:
            text_frame = body.text_frame
            text_frame.clear()
            for i, line in enumerate(content_lines):
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                p.text = line
                p.font.size = Pt(16)

    def _render_generic(
        self,
        root: etree._Element,
        prs: Presentation,
        citations: List[str],
    ) -> None:
        """Render generic XML element."""
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        title.text = root.tag.replace("-", " ").replace("_", " ").title()

        # Content slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = "Content"

        # Convert XML to text
        content = self._element_to_text(root, recursive=True)
        text_frame = body.text_frame
        text_frame.text = content

    def _element_to_text(
        self,
        element: etree._Element,
        recursive: bool = False,
    ) -> str:
        """Convert XML element to text."""
        parts = []

        if element.text and element.text.strip():
            parts.append(element.text.strip())

        if recursive:
            for child in element:
                child_text = self._element_to_text(child, recursive=True)
                if child_text:
                    parts.append(f"{child.tag}: {child_text}")

        return " ".join(parts)

    def _add_citations_slide(
        self,
        prs: Presentation,
        citations: List[str],
    ) -> None:
        """Add citations slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = "Citations & References"

        text_frame = body.text_frame
        text_frame.clear()

        for i, citation in enumerate(citations):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = f"[{i+1}] {citation}"
            p.font.size = Pt(10)
